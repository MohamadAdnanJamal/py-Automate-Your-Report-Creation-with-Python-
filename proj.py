import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from tkinter import filedialog, Tk, Button, Label, messagebox
import os
import sys
import ctypes
import subprocess

# Check if the script is running in a bundled application or in a development environment
if getattr(sys, 'frozen', False):
    # Running in a bundled application
    SCRIPT_DIR = sys._MEIPASS
else:
    # Running in a development environment
    SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# Function to create the PowerPoint presentation
def create_powerpoint():
    # Initialize Tkinter
    root = Tk()
    root.withdraw()  # Hide the root window

    # Ask the user to select the Excel file
    excel_file_path = filedialog.askopenfilename(title="Select Excel file", filetypes=[("Excel files", "*.xlsx")])
    if not excel_file_path:
        return  # Exit if no file selected

    # Read the Excel file
    df = pd.read_excel(excel_file_path)

    # Ask the user to select the location to save the new PowerPoint presentation
    pptx_save_path = filedialog.asksaveasfilename(title="Save PowerPoint file", defaultextension=".pptx",
                                                  filetypes=[("PowerPoint files", "*.pptx")])
    if not pptx_save_path:
        return  # Exit if no file selected

    # Create a presentation object
    prs = Presentation()

    # Define font settings
    font_calibri = "Calibri"
    font_size = Pt(18)  # Font size 18 points

    # Loop through each row in the DataFrame and create a slide for each row
    for index, row in df.iterrows():
        # Create a slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Remove default placeholder shapes
        for shape in slide.placeholders:
            if shape.placeholder_format.idx in [0, 1, 2, 3, 4]:  # Remove the first five placeholders
                sp = shape._sp
                sp.getparent().remove(sp)

        # Set up text box for columns A to D
        left_a_to_d = Inches(0.1)
        top_a_to_d = Inches(0.1)
        width_a_to_d = Inches(2)
        height_a_to_d = Inches(1.5)

        # Add a text box for columns A to D
        textbox_a_to_d = slide.shapes.add_textbox(left_a_to_d, top_a_to_d, width_a_to_d, height_a_to_d)
        text_frame_a_to_d = textbox_a_to_d.text_frame
        text_frame_a_to_d.word_wrap = True

        # Add text to the text box for columns A to D
        p_a_to_d = text_frame_a_to_d.add_paragraph()
        p_a_to_d.text = "\n".join(str(cell) if not pd.isna(cell) else "" for cell in row[:4])

        # Apply font settings to text in the text box for columns A to D
        for paragraph in text_frame_a_to_d.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_calibri
                run.font.size = font_size
            paragraph.alignment = PP_ALIGN.LEFT  # Left align the paragraph

        # Set up text box for column E
        left_e = Inches(3)
        top_e = Inches(6.5)
        width_e = Inches(2)
        height_e = Inches(0.7)

        # Add a text box for column E
        textbox_e = slide.shapes.add_textbox(left_e, top_e, width_e, height_e)
        text_frame_e = textbox_e.text_frame
        text_frame_e.word_wrap = True

        # Add text to the text box for column E
        p_e = text_frame_e.add_paragraph()
        p_e.text = str(row['E']) if not pd.isna(row['E']) else ""

        # Apply font settings to text in the text box for column E
        for paragraph in text_frame_e.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_calibri
                run.font.size = font_size
            paragraph.alignment = PP_ALIGN.CENTER  # Center align the paragraph

        # Set up text box for column F
        left_f = Inches(7)
        top_f = Inches(6.5)
        width_f = Inches(2)
        height_f = Inches(0.7)

        # Add a text box for column F
        textbox_f = slide.shapes.add_textbox(left_f, top_f, width_f, height_f)
        text_frame_f = textbox_f.text_frame
        text_frame_f.word_wrap = True

        # Add text to the text box for column F
        p_f = text_frame_f.add_paragraph()
        p_f.text = str(row['F']) if not pd.isna(row['F']) else ""

        # Apply font settings to text in the text box for column F
        for paragraph in text_frame_f.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_calibri
                run.font.size = font_size
            paragraph.alignment = PP_ALIGN.CENTER  # Center align the paragraph

        # Set up text box for column G
        left_g = Inches(0.1)
        top_g = Inches(2)
        width_g = Inches(2)
        height_g = Inches(0.7)

        # Add a text box for column G
        textbox_g = slide.shapes.add_textbox(left_g, top_g, width_g, height_g)
        text_frame_g = textbox_g.text_frame
        text_frame_g.word_wrap = True

        # Add text to the text box for column G
        p_g = text_frame_g.add_paragraph()
        p_g.text = str(row['G']) if not pd.isna(row['G']) else ""

        # Apply font settings to text in the text box for column G
        for paragraph in text_frame_g.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_calibri
                run.font.size = font_size
            paragraph.alignment = PP_ALIGN.LEFT  # Left align the paragraph

        # Set up text box for column H
        left_h = Inches(0.1)
        top_h = Inches(4)
        width_h = Inches(2)
        height_h = Inches(0.7)

        # Add a text box for column H
        textbox_h = slide.shapes.add_textbox(left_h, top_h, width_h, height_h)
        text_frame_h = textbox_h.text_frame
        text_frame_h.word_wrap = True

        # Add text to the text box for column H
        p_h = text_frame_h.add_paragraph()
        p_h.text = str(row['H']) if not pd.isna(row['H']) else ""

        # Apply font settings to text in the text box for column H
        for paragraph in text_frame_h.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_calibri
                run.font.size = font_size
            paragraph.alignment = PP_ALIGN.LEFT  # Left align the paragraph

    # Save the presentation
    prs.save(pptx_save_path)

    messagebox.showinfo("Success", "PowerPoint presentation saved successfully at:\n{}".format(pptx_save_path))

# Create a Tkinter window
window = Tk()
window.title("Excel to PowerPoint Converter")

# Set window size
window.geometry("300x150")

# Label
label = Label(window, text="Click below to select the Excel file and save the PowerPoint presentation:")
label.pack(pady=10)

# Button to create PowerPoint
create_button = Button(window, text="Create PowerPoint", command=create_powerpoint)
create_button.pack(pady=10)

# Run the Tkinter event loop
window.mainloop()
